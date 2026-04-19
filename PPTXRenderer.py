"""
BASF VIANT 动态 PPT 渲染器 - PPTXRenderer
基于 python-pptx 的动态渲染类，从 JSON Schema 生成专业 PPT
"""

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import nsmap
from enum import Enum
from typing import List, Optional, Dict, Any
from pydantic import BaseModel, Field, validator
import json
import os
import math


# ===== Pydantic 数据结构（用于规范 LLM 输出） =====

class SlideType(str, Enum):
    """幻灯片类型枚举"""
    COVER = "cover"                # 封面页
    TEXT_ONLY = "text_only"        # 纯文本页
    CHART_TEXT = "chart_text"      # 图表+文本页
    CHART_ONLY = "chart_only"      # 纯图表页

class ChartSeries(BaseModel):
    """图表数据系列"""
    name: str
    data: List[float]

class ChartData(BaseModel):
    """图表数据结构"""
    categories: List[str]
    series: List[ChartSeries]
    chart_type: str  # 'pie', 'bar', 'line', 'column'

class BulletPoint(BaseModel):
    """要点结构"""
    icon: str  # emoji 图标
    text: str  # 提炼过的句子

class Slide(BaseModel):
    """幻灯片结构"""
    slide_type: SlideType
    title: str
    subtitle: Optional[str] = None
    chart_data: Optional[ChartData] = None
    bullet_points: Optional[List[BulletPoint]] = None
    image_keyword: Optional[str] = None  # 用于图片搜索的英文关键词
    
    @validator('chart_data')
    def validate_chart_data(cls, v, values):
        slide_type = values.get('slide_type')
        if slide_type in [SlideType.CHART_TEXT, SlideType.CHART_ONLY] and v is None:
            raise ValueError(f'{slide_type} 类型的幻灯片必须包含 chart_data')
        if slide_type in [SlideType.COVER, SlideType.TEXT_ONLY] and v is not None:
            raise ValueError(f'{slide_type} 类型的幻灯片不应包含 chart_data')
        return v
    
    @validator('bullet_points')
    def validate_bullet_points(cls, v, values):
        slide_type = values.get('slide_type')
        if slide_type in [SlideType.TEXT_ONLY, SlideType.CHART_TEXT] and (v is None or len(v) == 0):
            raise ValueError(f'{slide_type} 类型的幻灯片必须包含 bullet_points')
        if slide_type in [SlideType.COVER, SlideType.CHART_ONLY] and v is not None:
            raise ValueError(f'{slide_type} 类型的幻灯片不应包含 bullet_points')
        return v
    
    @validator('image_keyword')
    def validate_image_keyword(cls, v, values):
        """验证图片关键词：如果提供，必须是英文"""
        if v is not None:
            # 确保是字符串且不包含中文字符
            import re
            if not isinstance(v, str):
                raise ValueError('image_keyword 必须是字符串')
            # 检查是否包含中文字符（简单的正则检查）
            if re.search(r'[\u4e00-\u9fff]', v):
                raise ValueError('image_keyword 必须是英文关键词，不包含中文')
        return v

class Presentation(BaseModel):
    """演示文稿结构"""
    title: str
    subtitle: Optional[str] = None
    date: str  # 格式: "August 2025"
    slides: List[Slide]


# ===== LLM 系统提示词 =====

LLM_SYSTEM_PROMPT = """
你是一位资深商业分析师，擅长从商业长文本中提取关键信息并转化为结构化数据。

## 任务
根据用户提供的商业长文本，提取关键信息并生成一个符合 JSON Schema 的演示文稿数据结构。

## 输出要求
你必须输出严格符合以下 Pydantic 模型的 JSON 数据：

```json
{
  "title": "string",
  "subtitle": "optional string",
  "date": "string",
  "slides": [
    {
      "slide_type": "cover | text_only | chart_text | chart_only",
      "title": "string",
      "subtitle": "optional string",
      "chart_data": {
        "categories": ["string"],
        "series": [
          {"name": "string", "data": [number]}
        ],
        "chart_type": "pie | bar | line | column"
      },
      "bullet_points": [
        {"icon": "emoji", "text": "string"}
      ],
      "image_keyword": "optional string"
    }
  ]
}
```

## 处理步骤
1. **提取元数据**：从文本中识别演示文稿的标题、副标题和日期。
2. **分析内容结构**：将文本内容分解为逻辑幻灯片，每张幻灯片应有明确的主题。
3. **确定幻灯片类型**：
   - 封面页 (cover)：包含标题、副标题、日期
   - 纯文本页 (text_only)：只有要点列表，无图表
   - 图表+文本页 (chart_text)：既有图表数据又有要点分析
   - 纯图表页 (chart_only)：只有图表，无要点文本
4. **提取图表数据**：
   - 从文本中识别数值数据，确保数据准确
   - 为图表选择合适的类型：饼图(pie)用于比例，柱状图(bar/column)用于比较，折线图(line)用于趋势
   - 确保 categories 和 series 数据对应正确
5. **提炼要点**：
   - 为每个核心观点提炼一句简洁的陈述
   - 为每个要点分配一个相关的 emoji 图标（如 🎯、📈、💰、🌍 等）
   - 要点应体现商业洞察力，而不仅仅是事实陈述
6. **判断是否需要插图**：
   - 如果幻灯片主题涉及以下内容，强烈建议提供 image_keyword：
     * Sustainability（可持续发展、环保、绿色技术）
     * Performance（性能、效率、数据展示）
     * Efficiency（效率、优化、流程改进）
     * 产品展示（chemical products, coatings, materials）
     * 技术概念（innovation, research, development）
     * 行业场景（factory, manufacturing, industrial）
   - 选择具体、有视觉表现力的英文关键词：
     * 好例子：sustainability, renewable energy, green chemistry, factory automation
     * 差例子：good, important, concept（过于抽象）
   - 封面页和纯图表页通常不需要插图

## 质量要求
- **商业洞察力**：要点应反映深度的商业分析，指出机会、风险、趋势等
- **数据准确**：图表数据必须与文本中的数字完全一致
- **结构合理**：幻灯片顺序应有逻辑，从概览到细节
- **简洁精炼**：每个要点不超过 15 个字，图表标题清晰明了
- **视觉增强**：在适当的幻灯片中提供 image_keyword 以增强视觉效果

## 示例
输入文本：
"2025年全球VIANT战略报告显示，我们在16个行业分析了800家制造商，识别出244家A级高机会客户。在Sustainability方面，VIANT产品可减少30%的碳排放..."

输出 JSON：
{
  "title": "Global VIANT Reconnecting",
  "subtitle": "Strategy 2025 Review",
  "date": "August 2025",
  "slides": [
    {
      "slide_type": "cover",
      "title": "Global VIANT Reconnecting",
      "subtitle": "Strategy 2025 Review"
    },
    {
      "slide_type": "chart_text",
      "title": "Executive Summary",
      "subtitle": "核心数据概览",
      "chart_data": {
        "categories": ["行业", "公司", "A级", "高机会"],
        "series": [
          {"name": "数量", "data": [16, 800, 244, 505]}
        ],
        "chart_type": "column"
      },
      "bullet_points": [
        {"icon": "🎯", "text": "16个行业深度分析，800家制造商覆盖"},
        {"icon": "⭐", "text": "244家A级高机会客户识别"},
        {"icon": "📈", "text": "平均匹配分数7.73分，数据质量优秀"}
      ]
    },
    {
      "slide_type": "text_only",
      "title": "Sustainability Impact",
      "subtitle": "环境效益分析",
      "bullet_points": [
        {"icon": "🌍", "text": "VIANT产品减少30%碳排放"},
        {"icon": "💧", "text": "节水40%的生产工艺"},
        {"icon": "♻️", "text": "95%材料可回收利用"}
      ],
      "image_keyword": "sustainable factory"
    }
  ]
}

现在，请根据用户提供的商业长文本，生成符合上述要求的结构化 JSON 数据。
"""


# ===== 动态 PPT 渲染器 =====

class PPTXRenderer:
    """动态 PPT 渲染器，根据 JSON Schema 生成专业 PPT"""
    
    def __init__(self, template_path='/root/owl/Template.pptx'):
        """
        初始化渲染器
        
        Args:
            template_path: PowerPoint 模板文件路径
        """
        self.template_path = template_path
        
        # BASF 品牌颜色
        self.BASF_BLUE = RGBColor(0, 44, 95)
        self.BASF_BLUE_LIGHT = RGBColor(20, 80, 140)  # 浅蓝变体
        self.BASF_GREEN = RGBColor(109, 185, 72)
        self.BASF_GREEN_LIGHT = RGBColor(140, 210, 120)
        self.BASF_LIGHT_BLUE = RGBColor(0, 159, 227)
        self.BASF_ORANGE = RGBColor(255, 140, 0)
        self.DARK_GREY = RGBColor(64, 64, 64)
        self.LIGHT_GREY = RGBColor(240, 240, 240)
        self.WHITE = RGBColor(255, 255, 255)
        self.MID_GREY = RGBColor(180, 180, 180)
        
        # 颜色池（用于图表系列着色）
        self.COLOR_POOL = [
            self.BASF_BLUE,
            self.BASF_LIGHT_BLUE,
            self.BASF_GREEN,
            self.BASF_ORANGE,
            self.BASF_BLUE_LIGHT,
            self.BASF_GREEN_LIGHT
        ]
        
        # 页面尺寸（16:9）- 精确尺寸 33.87cm x 19.05cm
        self.SLIDE_WIDTH = Inches(13.33)  # 33.87cm
        self.SLIDE_HEIGHT = Inches(7.5)   # 19.05cm
        
        # ===== 精确布局参数 =====
        # 表格区域（左侧）
        self.CHART_TOP = Inches(1.27)     # 距离上边框 3.24cm（原4cm，向上移动7.6mm）
        self.CHART_BOTTOM = Inches(0.94)  # 距离下边框 2.4cm
        self.CHART_LEFT = Inches(0.28)    # 距离左边框 0.7cm
        self.CHART_HEIGHT = Inches(5.12)  # 高 13cm
        self.CHART_WIDTH = Inches(6.3)    # 宽 16cm
        
        # 文本区域（右侧）
        self.TEXT_TOP = Inches(1.27)      # 距离上边框 3.24cm（原4cm，向上移动7.6mm）
        self.TEXT_BOTTOM = Inches(0.94)   # 距离下边框 2.4cm
        self.TEXT_RIGHT = Inches(0.28)    # 距离右边框 0.7cm
        self.TEXT_HEIGHT = Inches(5.12)   # 高 13cm
        self.TEXT_WIDTH = Inches(6.3)     # 宽 16cm
        
        # 计算文本区域 X 坐标（从右侧计算）
        self.TEXT_X = self.SLIDE_WIDTH - self.TEXT_RIGHT - self.TEXT_WIDTH
        
        # 页脚区域
        self.FOOTER_BOTTOM = Inches(0.2)  # 距离页面底部 0.5cm
        self.FOOTER_LEFT = Inches(0.28)   # 左边框 0.7cm
        self.FOOTER_HEIGHT = Inches(0.2)  # 高 0.5cm
        self.FOOTER_WIDTH = Inches(9.45)  # 长 24cm
        self.FOOTER_Y = self.SLIDE_HEIGHT - self.FOOTER_BOTTOM - self.FOOTER_HEIGHT
        
        # 表格和文本框中间的分隔线（调小宽度）
        self.SEPARATOR_WIDTH = Pt(1)      # 调小宽度
        
        # 图片布局参数
        self.IMAGE_TOP = Inches(1.27)      # 图片顶部与图表/文本对齐
        self.IMAGE_HEIGHT = Inches(3.0)    # 图片高度
        self.IMAGE_WIDTH = Inches(5.0)     # 图片宽度
        
        # 图片位置配置
        # 1. 左侧图片（当没有图表时）
        self.IMAGE_LEFT_X = Inches(0.28)   # 左边框 0.7cm
        # 2. 右侧图片（当有图表时，放在图表右侧）
        self.IMAGE_RIGHT_X = self.TEXT_X   # 与文本区域对齐
        # 3. 图表上方图片（较小的装饰性图片）
        self.IMAGE_SMALL_HEIGHT = Inches(1.5)
        self.IMAGE_SMALL_WIDTH = Inches(2.5)
        
        # 图片叠加效果参数
        self.IMAGE_OVERLAY_OPACITY = 0.2   # 叠加层透明度
        self.IMAGE_BORDER_WIDTH = Pt(2)    # 图片边框宽度
        self.IMAGE_BORDER_COLOR = self.BASF_BLUE_LIGHT  # 图片边框颜色
        
        # 图表类型映射
        self.CHART_TYPE_MAP = {
            'pie': XL_CHART_TYPE.PIE,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED
        }
    
    # ===== 公共渲染方法 =====
    
    def render_from_json(self, json_data: dict, output_path: str) -> str:
        """
        从 JSON 数据渲染 PPT
        
        Args:
            json_data: 符合 Presentation Schema 的 JSON 数据
            output_path: 输出 PPT 文件路径
            
        Returns:
            生成的 PPT 文件路径
        """
        # 验证 JSON 数据
        presentation = Presentation(**json_data)
        
        # 创建演示文稿
        prs = PPTXPresentation(self.template_path)
        
        # 渲染每一页
        for i, slide_data in enumerate(presentation.slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._render_slide(slide, slide_data, i + 1, presentation)
        
        # 保存文件
        prs.save(output_path)
        print(f"✅ PPT 已生成：{output_path}")
        print(f"📊 幻灯片数量：{len(presentation.slides)}")
        
        return output_path
    
    def render_from_presentation(self, presentation: Presentation, output_path: str) -> str:
        """
        从 Presentation 对象渲染 PPT
        
        Args:
            presentation: Presentation 对象
            output_path: 输出 PPT 文件路径
            
        Returns:
            生成的 PPT 文件路径
        """
        prs = PPTXPresentation(self.template_path)
        
        # 渲染每一页
        for i, slide_data in enumerate(presentation.slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._render_slide(slide, slide_data, i + 1, presentation)
        
        prs.save(output_path)
        print(f"✅ PPT 已生成：{output_path}")
        print(f"📊 幻灯片数量：{len(presentation.slides)}")
        
        return output_path
    
    # ===== 私有渲染方法 =====
    
    def _render_slide(self, slide, slide_data: Slide, page_num: int, presentation: Presentation):
        """
        渲染单张幻灯片
        
        Args:
            slide: pptx.slide.Slide 对象
            slide_data: Slide 数据对象
            page_num: 页码（从 1 开始）
            presentation: Presentation 对象（用于获取全局信息如日期）
        """
        # 根据幻灯片类型调用对应的渲染方法
        if slide_data.slide_type == SlideType.COVER:
            self._render_cover(slide, slide_data, page_num, presentation)
        elif slide_data.slide_type == SlideType.TEXT_ONLY:
            self._render_text_only(slide, slide_data, page_num)
        elif slide_data.slide_type == SlideType.CHART_TEXT:
            self._render_chart_text(slide, slide_data, page_num)
        elif slide_data.slide_type == SlideType.CHART_ONLY:
            self._render_chart_only(slide, slide_data, page_num)
        else:
            raise ValueError(f"未知的幻灯片类型：{slide_data.slide_type}")
        
        # 添加图片（如果有关键词）
        if slide_data.image_keyword:
            import os
            self._add_image_to_slide(slide, slide_data.image_keyword, slide_data)
    
    def _render_cover(self, slide, slide_data: Slide, page_num: int, presentation: Presentation):
        """渲染封面页"""
        # 使用模板中的标题占位符
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = slide_data.title
            title_tf = title.text_frame
            title_p = title_tf.paragraphs[0]
            title_p.font.size = Pt(38)
            title_p.font.bold = True
            title_p.font.name = 'Arial'
            title_p.font.color.rgb = self.WHITE
            title_p.alignment = PP_ALIGN.CENTER
        
        # 副标题（如果提供）
        if slide_data.subtitle:
            sub = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(0.6)
            )
            sub_tf = sub.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = slide_data.subtitle
            sub_p.font.size = Pt(30)
            sub_p.font.bold = True
            sub_p.font.name = 'Arial'
            sub_p.font.color.rgb = self.BASF_GREEN
            sub_p.alignment = PP_ALIGN.CENTER
        
        # 日期
        date = slide.shapes.add_textbox(
            Inches(1), Inches(4.5), Inches(8), Inches(0.4)
        )
        date_tf = date.text_frame
        date_p = date_tf.paragraphs[0]
        date_p.text = f"{presentation.date} | BASF Chemetall"
        date_p.font.size = Pt(15)
        date_p.font.name = 'Arial'
        date_p.font.color.rgb = self.WHITE
        date_p.alignment = PP_ALIGN.CENTER
        
        # 添加页脚
        self._add_footer(slide, page_num)
    
    def _render_text_only(self, slide, slide_data: Slide, page_num: int):
        """渲染纯文本页"""
        # 添加标题
        self._add_title(slide, slide_data.title, slide_data.subtitle)
        
        # 渲染要点卡片
        bullets = [(bp.icon, bp.text) for bp in slide_data.bullet_points]
        self._add_card_with_bullets(
            slide, 
            self.TEXT_X, 
            self.TEXT_TOP, 
            self.TEXT_WIDTH, 
            self.TEXT_HEIGHT, 
            bullets,
            title=slide_data.subtitle or "Key Points"
        )
        
        # 添加页脚
        self._add_footer(slide, page_num)
    
    def _render_chart_text(self, slide, slide_data: Slide, page_num: int):
        """渲染图表+文本页"""
        # 添加标题
        self._add_title(slide, slide_data.title, slide_data.subtitle)
        
        # 转换图表数据
        chart_data = self._convert_to_pptx_chart_data(slide_data.chart_data)
        chart_type = self.CHART_TYPE_MAP.get(
            slide_data.chart_data.chart_type, 
            XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        
        # 渲染图表
        self._add_chart_enhanced(
            slide,
            page_num,
            slide_data.title,
            slide_data.subtitle,
            chart_type,
            chart_data,
            slide_data.subtitle or "Chart",
            slide_data.bullet_points
        )
    
    def _render_chart_only(self, slide, slide_data: Slide, page_num: int):
        """渲染纯图表页"""
        # 添加标题
        self._add_title(slide, slide_data.title, slide_data.subtitle)
        
        # 转换图表数据
        chart_data = self._convert_to_pptx_chart_data(slide_data.chart_data)
        chart_type = self.CHART_TYPE_MAP.get(
            slide_data.chart_data.chart_type, 
            XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        
        # 渲染纯图表（无右侧文本区域）
        self._add_chart_only(
            slide,
            page_num,
            slide_data.title,
            slide_data.subtitle,
            chart_type,
            chart_data,
            slide_data.subtitle or "Chart"
        )
    
    # ===== 辅助渲染方法 =====
    
    def _add_title(self, slide, title_text, subtitle_text=None):
        """添加标题 - 使用模板中的标题占位符"""
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = title_text
            
            # 清理占位符中的 {{subtitle}} 等字样
            title_tf = title.text_frame
            title_tf.clear()  # 清除所有段落
            
            # 添加主标题
            title_p = title_tf.add_paragraph()
            title_p.text = title_text
            title_p.font.size = Pt(28)
            title_p.font.bold = True
            title_p.font.name = 'Arial'
            title_p.font.color.rgb = self.BASF_BLUE
            title_p.alignment = PP_ALIGN.CENTER
            
            # 添加副标题（如果有）
            if subtitle_text:
                subtitle_p = title_tf.add_paragraph()
                subtitle_p.text = subtitle_text
                subtitle_p.font.size = Pt(18)
                subtitle_p.font.name = 'Arial'
                subtitle_p.font.color.rgb = self.BASF_GREEN
                subtitle_p.alignment = PP_ALIGN.CENTER
                subtitle_p.space_before = Pt(8)
    
    def _add_chart_enhanced(self, slide, page_num, title, subtitle, chart_type, 
                           chart_data, chart_title, bullet_points=None):
        """
        添加增强版图表 - 使用精确布局参数
        """
        # 添加标题
        self._add_title(slide, title, subtitle)
        
        # 左侧图表背景卡片（圆角矩形）
        chart_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.CHART_LEFT - Inches(0.1), 
            self.CHART_TOP - Inches(0.05),
            self.CHART_WIDTH + Inches(0.2), 
            self.CHART_HEIGHT + Inches(0.1)
        )
        chart_bg.fill.solid()
        chart_bg.fill.fore_color.rgb = self.WHITE
        chart_bg.line.color.rgb = self.LIGHT_GREY
        chart_bg.line.width = Pt(1)
        
        # 图表（精确位置）
        chart_shape = slide.shapes.add_chart(
            chart_type, 
            self.CHART_LEFT, 
            self.CHART_TOP, 
            self.CHART_WIDTH, 
            self.CHART_HEIGHT, 
            chart_data
        )
        chart = chart_shape.chart
        
        # 图表标题优化
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.name = 'Arial'
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.BASF_BLUE
        
        # 自定义图表颜色（使用 BASF 品牌色池）
        if hasattr(chart.plots[0], 'series'):
            for i, series in enumerate(chart.plots[0].series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self.COLOR_POOL[i % len(self.COLOR_POOL)]
        
        # 右侧文字卡片（如果有要点）
        if bullet_points:
            bullets = [(bp.icon, bp.text) for bp in bullet_points]
            self._add_card_with_bullets(
                slide, 
                self.TEXT_X, 
                self.TEXT_TOP, 
                self.TEXT_WIDTH, 
                self.TEXT_HEIGHT, 
                bullets,
                title=subtitle or "Key Insights"
            )
        
        # 页脚
        self._add_footer(slide, page_num)
    
    def _add_chart_only(self, slide, page_num, title, subtitle, chart_type, 
                       chart_data, chart_title):
        """
        添加纯图表（无右侧文本区域）
        """
        # 添加标题
        self._add_title(slide, title, subtitle)
        
        # 图表背景卡片（居中，占用更多空间）
        chart_bg_width = self.SLIDE_WIDTH - Inches(1.5)
        chart_bg_left = (self.SLIDE_WIDTH - chart_bg_width) / 2
        chart_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            chart_bg_left, 
            self.CHART_TOP - Inches(0.05),
            chart_bg_width, 
            self.CHART_HEIGHT + Inches(0.1)
        )
        chart_bg.fill.solid()
        chart_bg.fill.fore_color.rgb = self.WHITE
        chart_bg.line.color.rgb = self.LIGHT_GREY
        chart_bg.line.width = Pt(1)
        
        # 图表（居中放大）
        chart_width = chart_bg_width - Inches(0.5)
        chart_left = chart_bg_left + Inches(0.25)
        chart_shape = slide.shapes.add_chart(
            chart_type, 
            chart_left, 
            self.CHART_TOP, 
            chart_width, 
            self.CHART_HEIGHT, 
            chart_data
        )
        chart = chart_shape.chart
        
        # 图表标题优化
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.name = 'Arial'
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.BASF_BLUE
        
        # 自定义图表颜色
        if hasattr(chart.plots[0], 'series'):
            for i, series in enumerate(chart.plots[0].series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self.COLOR_POOL[i % len(self.COLOR_POOL)]
        
        # 页脚
        self._add_footer(slide, page_num)
    
    def _add_card_with_bullets(self, slide, x, y, w, h, bullets, title=None, icon='📊'):
        """
        添加卡片式 Bullet Points 区域（自适应版本）
        
        Args:
            bullets: 列表，每个元素为 (icon, text) 元组
        """
        # 卡片背景（圆角矩形）
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, w, h
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.WHITE
        bg.line.color.rgb = self.LIGHT_GREY
        bg.line.width = Pt(1)
        
        # 添加左侧色条装饰
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y, Inches(0.25), h
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.BASF_BLUE
        accent.line.fill.background()
        
        # 内容文本框
        tf = slide.shapes.add_textbox(
            x + Inches(0.35), 
            y + Inches(0.25),
            w - Inches(0.5), 
            h - Inches(0.5)
        ).text_frame
        tf.word_wrap = True
        tf.clear()
        
        # 动态计算字体大小和行距
        base_font_size = self._calculate_adaptive_font_size(len(bullets), h)
        line_spacing = self._calculate_line_spacing(len(bullets))
        
        # 标题（如果有）
        if title:
            title_p = tf.add_paragraph()
            title_p.text = title
            title_p.font.size = Pt(base_font_size + 2)  # 标题稍大
            title_p.font.bold = True
            title_p.font.name = 'Arial'
            title_p.font.color.rgb = self.BASF_BLUE
            title_p.space_after = Pt(12)
        
        # Bullet Points
        for i, (icon_char, text) in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = f"{icon_char}  {text}"
            p.font.size = Pt(base_font_size)
            p.font.name = 'Arial'
            p.font.color.rgb = self.DARK_GREY
            p.line_spacing = line_spacing
            
            if i > 0:
                p.space_before = Pt(6)
    
    def _add_footer(self, slide, page_num):
        """添加页脚（移除蓝色分割线）"""
        # 页脚背景（浅灰色）
        footer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.FOOTER_LEFT, self.FOOTER_Y,
            self.FOOTER_WIDTH, self.FOOTER_HEIGHT
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = self.LIGHT_GREY
        footer.line.fill.background()
        
        # 左侧文本：BASF Chemetall
        left_text = slide.shapes.add_textbox(
            self.FOOTER_LEFT + Inches(0.1), 
            self.FOOTER_Y,
            Inches(3), 
            self.FOOTER_HEIGHT
        ).text_frame
        left_p = left_text.add_paragraph()
        left_p.text = "BASF Chemetall - VIANT Reconnecting"
        left_p.font.size = Pt(9)
        left_p.font.name = 'Arial'
        left_p.font.color.rgb = self.MID_GREY
        
        # 右侧页码
        right_text = slide.shapes.add_textbox(
            self.FOOTER_LEFT + self.FOOTER_WIDTH - Inches(1), 
            self.FOOTER_Y,
            Inches(0.8), 
            self.FOOTER_HEIGHT
        ).text_frame
        right_p = right_text.add_paragraph()
        right_p.text = str(page_num)
        right_p.font.size = Pt(9)
        right_p.font.bold = True
        right_p.font.name = 'Arial'
        right_p.font.color.rgb = self.MID_GREY
        right_p.alignment = PP_ALIGN.RIGHT
    
    # ===== 工具方法 =====
    
    def _convert_to_pptx_chart_data(self, chart_data: ChartData) -> CategoryChartData:
        """
        将 ChartData 模型转换为 pptx CategoryChartData
        """
        pptx_data = CategoryChartData()
        pptx_data.categories = chart_data.categories
        
        for series in chart_data.series:
            pptx_data.add_series(series.name, series.data)
        
        return pptx_data
    
    def _calculate_adaptive_font_size(self, bullet_count: int, available_height: float) -> float:
        """
        根据要点数量和可用高度计算自适应字体大小
        
        Args:
            bullet_count: 要点数量
            available_height: 可用高度（英寸）
            
        Returns:
            建议的字体大小（磅）
        """
        # 基础字体大小
        base_size = 12.5
        
        # 根据数量调整
        if bullet_count <= 3:
            return base_size + 1  # 13.5
        elif bullet_count <= 5:
            return base_size  # 12.5
        elif bullet_count <= 8:
            return base_size - 1  # 11.5
        else:
            return max(base_size - 2, 10)  # 最小 10
    
    def _calculate_line_spacing(self, bullet_count: int) -> float:
        """
        根据要点数量计算行距
        
        Args:
            bullet_count: 要点数量
            
        Returns:
            行距倍数
        """
        if bullet_count <= 3:
            return 1.6  # 宽松行距
        elif bullet_count <= 5:
            return 1.4  # 适中行距
        elif bullet_count <= 8:
            return 1.2  # 紧凑行距
        else:
            return 1.1  # 非常紧凑
    
    def fetch_image_for_slide(self, keyword: str) -> str:
        """
        获取幻灯片配图（占位实现）
        
        Args:
            keyword: 图片搜索关键词（英文）
            
        Returns:
            图片文件路径
            
        注意：这是一个占位实现。在实际部署中，可以替换为：
        1. 调用 Unsplash/Pexels API 获取免费图片
        2. 使用本地图片库
        3. 调用商业图片API
        """
        # 记录图片请求（实际使用时替换为真实API调用）
        print(f"📷 请求图片关键词: {keyword}")
        
        # 创建测试图片目录
        import os
        test_image_dir = "/tmp/ppt_images"
        os.makedirs(test_image_dir, exist_ok=True)
        
        # 根据关键词生成不同的测试图片
        # 这里使用PIL创建简单的测试图片，如果没有PIL则使用纯色图片文件
        try:
            from PIL import Image, ImageDraw, ImageFont
            import hashlib
            
            # 根据关键词生成确定性颜色
            hash_val = int(hashlib.md5(keyword.encode()).hexdigest()[:6], 16)
            r = (hash_val >> 16) & 0xFF
            g = (hash_val >> 8) & 0xFF
            b = hash_val & 0xFF
            
            # 创建图片
            img_width, img_height = 800, 600
            image = Image.new('RGB', (img_width, img_height), color=(r, g, b))
            draw = ImageDraw.Draw(image)
            
            # 添加文字
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 40)
            except:
                font = ImageFont.load_default()
            
            # 在图片中心添加关键词
            text = f"{keyword}"
            # 计算文字位置
            from PIL import ImageFont
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            position = ((img_width - text_width) // 2, (img_height - text_height) // 2)
            
            # 绘制文字
            draw.text(position, text, fill=(255, 255, 255), font=font)
            
            # 保存图片
            safe_keyword = "".join(c for c in keyword if c.isalnum() or c in (' ', '-', '_')).rstrip()
            image_path = os.path.join(test_image_dir, f"{safe_keyword[:50]}.jpg")
            image.save(image_path, "JPEG", quality=95)
            
            print(f"   ✅ 生成测试图片: {image_path}")
            return image_path
            
        except ImportError:
            # 如果没有PIL，创建简单的文本文件作为占位
            safe_keyword = "".join(c for c in keyword if c.isalnum() or c in (' ', '-', '_')).rstrip()
            image_path = os.path.join(test_image_dir, f"{safe_keyword[:50]}.txt")
            
            with open(image_path, 'w') as f:
                f.write(f"Placeholder image for: {keyword}\n")
                f.write(f"Actual implementation would fetch from image API.\n")
            
            print(f"   ⚠️  PIL未安装，创建文本占位文件: {image_path}")
            return image_path
    
    def _add_image_to_slide(self, slide, image_keyword: str, slide_data: Slide):
        """
        向幻灯片添加图片
        
        Args:
            slide: pptx.slide.Slide 对象
            image_keyword: 图片关键词
            slide_data: Slide 数据对象
        """
        if not image_keyword:
            return
        
        try:
            # 获取图片
            image_path = self.fetch_image_for_slide(image_keyword)
            if not image_path or not os.path.exists(image_path):
                print(f"   ⚠️ 图片未找到: {image_keyword}")
                return
            
            # 根据幻灯片类型决定图片位置和大小
            if slide_data.slide_type == SlideType.COVER:
                # 封面页：不添加图片或添加背景图
                return
                
            elif slide_data.slide_type == SlideType.CHART_ONLY:
                # 纯图表页：在图表上方添加小图
                image_left = self.CHART_LEFT + (self.CHART_WIDTH - self.IMAGE_SMALL_WIDTH) / 2
                image_top = self.IMAGE_TOP - self.IMAGE_SMALL_HEIGHT - Inches(0.2)
                image_shape = slide.shapes.add_picture(
                    image_path,
                    image_left, image_top,
                    width=self.IMAGE_SMALL_WIDTH,
                    height=self.IMAGE_SMALL_HEIGHT
                )
                self._create_image_overlay(slide, image_shape)
                
            elif slide_data.slide_type == SlideType.CHART_TEXT:
                # 图表+文本页：在文本区域添加图片（右侧）
                image_left = self.IMAGE_RIGHT_X
                image_top = self.IMAGE_TOP
                image_shape = slide.shapes.add_picture(
                    image_path,
                    image_left, image_top,
                    width=self.IMAGE_WIDTH,
                    height=self.IMAGE_HEIGHT
                )
                self._create_image_overlay(slide, image_shape)
                
                # 调整文本区域位置（为图片腾出空间）
                # 文本区域向右移动，高度减少
                pass  # 可以在后续调整
                
            elif slide_data.slide_type == SlideType.TEXT_ONLY:
                # 纯文本页：在左侧添加大图
                image_left = self.IMAGE_LEFT_X
                image_top = self.IMAGE_TOP
                image_shape = slide.shapes.add_picture(
                    image_path,
                    image_left, image_top,
                    width=self.IMAGE_WIDTH,
                    height=self.IMAGE_HEIGHT
                )
                self._create_image_overlay(slide, image_shape)
                
                # 调整文本区域位置（为图片腾出空间）
                # 文本区域向右移动
                pass  # 可以在后续调整
            
            print(f"   ✅ 图片插入成功: {image_keyword}")
            
        except Exception as e:
            print(f"   ❌ 图片插入失败: {e}")
            import traceback
            traceback.print_exc()
    
    def _create_image_overlay(self, slide, image_shape):
        """
        创建图片叠加效果
        
        Args:
            slide: pptx.slide.Slide 对象
            image_shape: 图片形状对象
        """
        try:
            # 添加半透明叠加层
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                image_shape.left,
                image_shape.top,
                image_shape.width,
                image_shape.height
            )
            
            # 设置填充颜色和透明度
            fill = overlay.fill
            fill.solid()
            fill.fore_color.rgb = self.BASF_BLUE
            fill.transparency = self.IMAGE_OVERLAY_OPACITY
            
            # 设置边框
            line = overlay.line
            line.color.rgb = self.IMAGE_BORDER_COLOR
            line.width = self.IMAGE_BORDER_WIDTH
            line.dash_style = MSO_LINE_DASH_STYLE.SOLID
            
            # 将叠加层置于图片上方
            overlay.z_order(1)
            
            # 调整图片边框
            img_line = image_shape.line
            img_line.color.rgb = self.WHITE
            img_line.width = Pt(1)
            
        except Exception as e:
            print(f"   叠加效果创建失败: {e}")
    
    def get_system_prompt(self) -> str:
        """获取 LLM 系统提示词"""
        return LLM_SYSTEM_PROMPT


# ===== 使用示例 =====

def example_usage():
    """使用示例"""
    # 创建渲染器
    renderer = PPTXRenderer('/root/owl/Template.pptx')
    
    # 示例 JSON 数据
    example_json = {
        "title": "Global VIANT Reconnecting",
        "subtitle": "Strategy 2025 Review",
        "date": "August 2025",
        "slides": [
            {
                "slide_type": "cover",
                "title": "Global VIANT Reconnecting",
                "subtitle": "Strategy 2025 Review"
            },
            {
                "slide_type": "chart_text",
                "title": "Executive Summary",
                "subtitle": "核心数据概览",
                "chart_data": {
                    "categories": ["行业", "公司", "A级", "高机会"],
                    "series": [
                        {"name": "数量", "data": [16, 800, 244, 505]}
                    ],
                    "chart_type": "column"
                },
                "bullet_points": [
                    {"icon": "🎯", "text": "16个行业深度分析，800家制造商覆盖"},
                    {"icon": "⭐", "text": "244家A级高机会客户识别"},
                    {"icon": "📈", "text": "平均匹配分数7.73分，数据质量优秀"}
                ]
            },
            {
                "slide_type": "text_only",
                "title": "Strategic Insights",
                "subtitle": "关键发现",
                "bullet_points": [
                    {"icon": "💰", "text": "总机会规模 $591.5M"},
                    {"icon": "📈", "text": "预期收入 $146.5M (24.5% 转化率)"},
                    {"icon": "🌿", "text": "HDG 贡献 63% 预期收入"},
                    {"icon": "🎯", "text": "年度目标：签约 60 家 A 级"},
                    {"icon": "✅", "text": "成功概率评估：高 (85%)"}
                ],
                "image_keyword": "sustainable factory"
            },
            {
                "slide_type": "chart_text",
                "title": "Sustainability Impact",
                "subtitle": "环境效益分析",
                "chart_data": {
                    "categories": ["碳排放", "用水量", "材料回收"],
                    "series": [
                        {"name": "减少比例", "data": [30, 40, 95]}
                    ],
                    "chart_type": "column"
                },
                "bullet_points": [
                    {"icon": "🌍", "text": "VIANT产品减少30%碳排放"},
                    {"icon": "💧", "text": "节水40%的生产工艺"},
                    {"icon": "♻️", "text": "95%材料可回收利用"}
                ],
                "image_keyword": "green technology"
            }
        ]
    }
    
    # 渲染 PPT
    output_path = renderer.render_from_json(
        example_json,
        '/root/VIANT_Dynamic_Output.pptx'
    )
    
    print(f"🎉 示例 PPT 已生成：{output_path}")
    print(f"📋 使用提示词长度：{len(renderer.get_system_prompt())} 字符")


if __name__ == '__main__':
    example_usage()