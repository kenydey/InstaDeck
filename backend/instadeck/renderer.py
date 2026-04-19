"""PPTX rendering: native charts + text."""

from __future__ import annotations

import io
from pathlib import Path

from pptx import Presentation as PPTXPresentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from instadeck.schemas import ChartData, Presentation, Slide, SlideType

THEME_COLORS = [
    RGBColor(0, 82, 155),
    RGBColor(0, 159, 227),
    RGBColor(109, 185, 72),
    RGBColor(255, 140, 0),
    RGBColor(120, 120, 120),
]

CHART_MAP = {
    "pie": XL_CHART_TYPE.PIE,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "waterfall": getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED),
    "combo": XL_CHART_TYPE.COLUMN_CLUSTERED,
}


def _blank_layout(prs: PPTXPresentation):
    for idx in (6, 5, 1, 0):
        if idx < len(prs.slide_layouts):
            return prs.slide_layouts[idx]
    return prs.slide_layouts[0]


def _font_scale_for_bullets(n: int) -> float:
    if n <= 4:
        return 1.0
    if n <= 6:
        return 0.9
    return 0.78


def _add_chart(slide, chart_data: ChartData, x, y, cx, cy):
    chart_data_c = CategoryChartData()
    chart_data_c.categories = chart_data.categories
    for s in chart_data.series:
        chart_data_c.add_series(s.name, tuple(s.data))
    ctype = CHART_MAP.get(chart_data.chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    graphic_frame = slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data_c)
    chart = graphic_frame.chart
    for i, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = THEME_COLORS[i % len(THEME_COLORS)]


def _render_cover(slide, s: Slide, prs: Presentation):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = s.title
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if s.subtitle:
        p2 = tf.add_paragraph()
        p2.text = s.subtitle
        p2.font.size = Pt(22)
        p2.alignment = PP_ALIGN.CENTER
    foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    foot.text_frame.text = prs.date


def _render_text(slide, s: Slide):
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    t.text_frame.text = s.title
    t.text_frame.paragraphs[0].font.size = Pt(28)
    t.text_frame.paragraphs[0].font.bold = True
    n = len(s.bullet_points or [])
    fs = int(18 * _font_scale_for_bullets(n))
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(11.5), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bp in enumerate(s.bullet_points or []):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"{bp.icon}  {bp.text}".strip()
        para.font.size = Pt(fs)
        para.space_after = Pt(8)


def _render_chart_only(slide, s: Slide):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    title.text_frame.text = s.title
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    if s.chart_data:
        _add_chart(slide, s.chart_data, Inches(1), Inches(1.2), Inches(11), Inches(5.5))


def _render_chart_text(slide, s: Slide):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title.text_frame.text = s.title
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    if s.chart_data:
        _add_chart(slide, s.chart_data, Inches(0.4), Inches(1.0), Inches(6.2), Inches(5.0))
    n = len(s.bullet_points or [])
    fs = int(14 * _font_scale_for_bullets(n))
    tb = slide.shapes.add_textbox(Inches(6.9), Inches(1.0), Inches(6.0), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bp in enumerate(s.bullet_points or []):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"{bp.icon}  {bp.text}".strip()
        para.font.size = Pt(fs)


def render_presentation_to_path(
    presentation: Presentation,
    output_path: Path,
    template_path: Path | None = None,
) -> Path:
    """Build pptx from Presentation model."""
    if template_path and template_path.exists():
        prs = PPTXPresentation(str(template_path))
    else:
        prs = PPTXPresentation()
    layout = _blank_layout(prs)
    # reuse first slide if deck only has default one empty
    first = True
    for slide_model in presentation.slides:
        if first and len(prs.slides) > 0:
            slide = prs.slides[0]
            # clear placeholder shapes if any
            for shape in list(slide.shapes):
                sp = shape._element  # type: ignore[attr-defined]
                sp.getparent().remove(sp)
            first = False
        else:
            slide = prs.slides.add_slide(layout)
        if slide_model.slide_type == SlideType.cover:
            _render_cover(slide, slide_model, presentation)
        elif slide_model.slide_type == SlideType.text_only:
            _render_text(slide, slide_model)
        elif slide_model.slide_type == SlideType.chart_only:
            _render_chart_only(slide, slide_model)
        elif slide_model.slide_type == SlideType.chart_text:
            _render_chart_text(slide, slide_model)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    buf = io.BytesIO()
    prs.save(buf)
    output_path.write_bytes(buf.getvalue())
    return output_path
