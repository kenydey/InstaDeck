"""
端到端 PPT 生成管道
将文档解析、LLM 结构化提取和 PPT 渲染串联起来
"""

import os
import sys
import json
import argparse
from typing import Optional, Dict, Any
from datetime import datetime

# 导入自定义模块
from DocumentParser import DocumentParser
from LLMClient import LLMClient
from PPTXRenderer import PPTXRenderer


class DocumentToPPTPipeline:
    """文档到 PPT 的端到端管道"""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        初始化管道
        
        Args:
            config: 配置字典，包含各模块的配置参数
        """
        self.config = config or {}
        
        # 初始化各模块
        self.document_parser = DocumentParser(
            prefer_pdfplumber=self.config.get("prefer_pdfplumber", True)
        )
        
        # LLM 客户端配置
        llm_config = self.config.get("llm", {})
        self.llm_client = LLMClient(
            mode=llm_config.get("mode", "mock"),
            **llm_config.get("params", {})
        )
        
        # PPT 渲染器配置
        ppt_config = self.config.get("ppt", {})
        self.ppt_renderer = PPTXRenderer(
            template_path=ppt_config.get("template_path", "/root/owl/Template.pptx")
        )
    
    def run(self, input_path: str, output_path: Optional[str] = None) -> str:
        """
        运行完整管道
        
        Args:
            input_path: 输入文档路径
            output_path: 输出 PPT 路径（可选，默认自动生成）
            
        Returns:
            生成的 PPT 文件路径
        """
        print("=" * 60)
        print("🚀 开始文档到 PPT 转换管道")
        print(f"📄 输入文件: {input_path}")
        print("=" * 60)
        
        # 步骤 1: 解析文档
        print("\n📝 步骤 1: 解析文档...")
        try:
            document_text = self.document_parser.parse(input_path)
            print(f"   ✅ 文档解析完成，长度: {len(document_text)} 字符")
            
            # 可选：显示文档预览
            if self.config.get("show_document_preview", True):
                preview = document_text[:500] + "..." if len(document_text) > 500 else document_text
                print(f"   📋 文档预览:\n{preview}\n")
                
        except Exception as e:
            print(f"   ❌ 文档解析失败: {e}")
            raise
        
        # 步骤 2: LLM 提取结构化数据
        print("\n🤖 步骤 2: LLM 提取结构化数据...")
        try:
            llm_params = self.config.get("llm", {}).get("extract_params", {})
            presentation_data = self.llm_client.extract_presentation_data(document_text, **llm_params)
            print(f"   ✅ LLM 提取完成，幻灯片数量: {len(presentation_data.get('slides', []))}")
            
            # 可选：显示提取的数据预览
            if self.config.get("show_llm_output", False):
                print(f"   📊 提取的数据预览:")
                print(json.dumps(presentation_data, indent=2, ensure_ascii=False)[:1000] + "...")
                
        except Exception as e:
            print(f"   ❌ LLM 提取失败: {e}")
            raise
        
        # 步骤 3: 验证和修复数据
        print("\n🔍 步骤 3: 验证数据格式...")
        try:
            # LLMClient 已经进行了验证，这里进行双重检查
            from PPTXRenderer import Presentation
            presentation = Presentation(**presentation_data)
            print(f"   ✅ 数据验证通过，共 {len(presentation.slides)} 张幻灯片")
            
        except Exception as e:
            print(f"   ⚠️ 数据验证发现问题: {e}")
            print("   尝试修复数据...")
            try:
                presentation_data = self.llm_client.validate_and_fix_data(presentation_data)
                presentation = Presentation(**presentation_data)
                print(f"   ✅ 数据修复成功，共 {len(presentation.slides)} 张幻灯片")
            except Exception as e2:
                print(f"   ❌ 数据修复失败: {e2}")
                raise
        
        # 步骤 4: 生成 PPT
        print("\n🎨 步骤 4: 生成 PowerPoint 演示文稿...")
        
        # 自动生成输出路径（如果未提供）
        if not output_path:
            base_name = os.path.splitext(os.path.basename(input_path))[0]
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"/tmp/{base_name}_presentation_{timestamp}.pptx"
        
        try:
            # 使用渲染器生成 PPT
            result_path = self.ppt_renderer.render_from_json(presentation_data, output_path)
            print(f"   ✅ PPT 生成成功: {result_path}")
            
            # 显示文件信息
            if os.path.exists(result_path):
                size_mb = os.path.getsize(result_path) / (1024 * 1024)
                print(f"   📊 文件大小: {size_mb:.2f} MB")
            
            return result_path
            
        except Exception as e:
            print(f"   ❌ PPT 生成失败: {e}")
            raise
    
    def run_with_sections(self, input_path: str, output_path: Optional[str] = None) -> str:
        """
        使用文档章节分割运行管道（更精细的控制）
        
        Args:
            input_path: 输入文档路径
            output_path: 输出 PPT 路径
            
        Returns:
            生成的 PPT 文件路径
        """
        print("=" * 60)
        print("🚀 开始章节化文档到 PPT 转换管道")
        print("=" * 60)
        
        # 步骤 1: 解析文档并分割章节
        print("\n📝 步骤 1: 解析文档并分割章节...")
        try:
            sections = self.document_parser.parse_to_sections(input_path)
            print(f"   ✅ 文档解析完成，共 {len(sections)} 个章节")
            
            # 显示章节概览
            for i, (title, content) in enumerate(sections[:5]):  # 只显示前5个
                print(f"   {i+1}. {title} ({len(content)} 字符)")
            if len(sections) > 5:
                print(f"   ... 还有 {len(sections) - 5} 个章节")
                
            # 合并所有章节内容
            document_text = "\n\n".join([f"## {title}\n{content}" for title, content in sections])
            
        except Exception as e:
            print(f"   ❌ 文档解析失败: {e}")
            raise
        
        # 步骤 2-4: 与常规流程相同
        return self.run_from_text(document_text, output_path, f"基于章节的 {os.path.basename(input_path)}")
    
    def run_from_text(self, text: str, output_path: Optional[str] = None, title_hint: str = "") -> str:
        """
        直接从文本运行管道
        
        Args:
            text: 输入文本
            output_path: 输出 PPT 路径
            title_hint: 标题提示（用于生成输出文件名）
            
        Returns:
            生成的 PPT 文件路径
        """
        print(f"\n📝 直接处理文本，长度: {len(text)} 字符")
        
        # 自动生成输出路径（如果未提供）
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            title_part = title_hint[:20] if title_hint else "presentation"
            output_path = f"/tmp/{title_part}_{timestamp}.pptx"
        
        # 步骤 2: LLM 提取结构化数据
        print("\n🤖 步骤 2: LLM 提取结构化数据...")
        try:
            llm_params = self.config.get("llm", {}).get("extract_params", {})
            presentation_data = self.llm_client.extract_presentation_data(text, **llm_params)
            print(f"   ✅ LLM 提取完成，幻灯片数量: {len(presentation_data.get('slides', []))}")
            
        except Exception as e:
            print(f"   ❌ LLM 提取失败: {e}")
            raise
        
        # 步骤 3: 验证数据
        print("\n🔍 步骤 3: 验证数据格式...")
        try:
            from PPTXRenderer import Presentation
            presentation = Presentation(**presentation_data)
            print(f"   ✅ 数据验证通过，共 {len(presentation.slides)} 张幻灯片")
            
        except Exception as e:
            print(f"   ⚠️ 数据验证发现问题: {e}")
            print("   尝试修复数据...")
            try:
                presentation_data = self.llm_client.validate_and_fix_data(presentation_data)
                presentation = Presentation(**presentation_data)
                print(f"   ✅ 数据修复成功，共 {len(presentation.slides)} 张幻灯片")
            except Exception as e2:
                print(f"   ❌ 数据修复失败: {e2}")
                raise
        
        # 步骤 4: 生成 PPT
        print("\n🎨 步骤 4: 生成 PowerPoint 演示文稿...")
        try:
            result_path = self.ppt_renderer.render_from_json(presentation_data, output_path)
            print(f"   ✅ PPT 生成成功: {result_path}")
            
            if os.path.exists(result_path):
                size_mb = os.path.getsize(result_path) / (1024 * 1024)
                print(f"   📊 文件大小: {size_mb:.2f} MB")
            
            return result_path
            
        except Exception as e:
            print(f"   ❌ PPT 生成失败: {e}")
            raise


def create_default_config() -> Dict[str, Any]:
    """创建默认配置"""
    return {
        # 文档解析配置
        "prefer_pdfplumber": True,
        "show_document_preview": True,
        "show_llm_output": False,
        
        # LLM 配置
        "llm": {
            "mode": "mock",  # 默认使用模拟模式
            "params": {
                # API 参数（模拟模式下忽略）
                "model": "gpt-4-turbo-preview",
                "temperature": 0.2,
                "max_tokens": 4000
            },
            "extract_params": {
                # 提取时的特定参数
            }
        },
        
        # PPT 配置
        "ppt": {
            "template_path": "/root/owl/Template.pptx"
        }
    }


def main():
    """命令行入口点"""
    parser = argparse.ArgumentParser(description="从文档生成 PowerPoint 演示文稿")
    parser.add_argument("input", help="输入文档路径（支持 .txt, .md, .docx, .pdf）")
    parser.add_argument("-o", "--output", help="输出 PPT 路径（可选，默认自动生成）")
    parser.add_argument("--sections", action="store_true", 
                       help="使用章节分割模式（对长文档更有效）")
    parser.add_argument("--llm-mode", choices=["mock", "openai", "openrouter"], 
                       default="mock", help="LLM 模式（默认: mock）")
    parser.add_argument("--api-key", help="API 密钥（如果使用真实 LLM）")
    parser.add_argument("--base-url", help="API 基础 URL（如果使用真实 LLM）")
    parser.add_argument("--model", help="LLM 模型名称")
    parser.add_argument("--show-llm-output", action="store_true", 
                       help="显示 LLM 输出内容")
    parser.add_argument("--config", help="JSON 配置文件路径")
    
    args = parser.parse_args()
    
    # 加载配置
    if args.config and os.path.exists(args.config):
        with open(args.config, 'r', encoding='utf-8') as f:
            config = json.load(f)
    else:
        config = create_default_config()
    
    # 覆盖命令行参数
    config["llm"]["mode"] = args.llm_mode
    config["show_llm_output"] = args.show_llm_output
    
    if args.api_key:
        config["llm"]["params"]["api_key"] = args.api_key
    if args.base_url:
        config["llm"]["params"]["base_url"] = args.base_url
    if args.model:
        config["llm"]["params"]["model"] = args.model
    
    # 检查输入文件
    if not os.path.exists(args.input):
        print(f"❌ 错误：输入文件不存在: {args.input}")
        sys.exit(1)
    
    # 检查模板文件
    template_path = config["ppt"]["template_path"]
    if not os.path.exists(template_path):
        print(f"⚠️ 警告：模板文件不存在: {template_path}")
        print("   将尝试使用默认模板...")
        # 尝试寻找其他模板
        potential_paths = [
            "/root/owl/Template.pptx",
            "./owl/Template.pptx",
            "./Template.pptx"
        ]
        for path in potential_paths:
            if os.path.exists(path):
                config["ppt"]["template_path"] = path
                print(f"   使用模板: {path}")
                break
    
    # 创建并运行管道
    try:
        pipeline = DocumentToPPTPipeline(config)
        
        if args.sections:
            result_path = pipeline.run_with_sections(args.input, args.output)
        else:
            result_path = pipeline.run(args.input, args.output)
        
        print("\n" + "=" * 60)
        print(f"🎉 转换完成！")
        print(f"📁 输出文件: {result_path}")
        print("=" * 60)
        
        # 提示如何发送文件
        if "telegram" in sys.modules or "send_message" in globals():
            print("\n💡 提示：您可以使用 send_message 工具将此文件发送到 Telegram")
        
    except Exception as e:
        print(f"\n❌ 管道执行失败: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    # 检查必要的依赖
    try:
        import docx
        import pdfplumber
        import pydantic
        from pptx import Presentation
    except ImportError as e:
        print(f"❌ 缺少依赖: {e}")
        print("请安装必要的 Python 包：")
        print("  pip install python-docx pdfplumber pydantic python-pptx")
        print("或使用现有环境：")
        print("  source /tmp/ppt_env/bin/activate")
        sys.exit(1)
    
    main()